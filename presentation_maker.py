from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

slide1 = prs.slides.add_slide(prs.slide_layouts[5]) #adding a new slide
title1 = slide1.shapes.title
title1.text = "First Slide" # title of first slide

slide2 = prs.slides.add_slide(prs.slide_layouts[5])
title2 = slide2.shapes.title
title2.text = "Second slide"
content2 = slide2.shapes.placeholders[0].text_frame
content2.text = "This is my second slide welcome to second slide"

# Slide 3: Conclusion
slide3 = prs.slides.add_slide(prs.slide_layouts[5])
title3 = slide3.shapes.title
title3.text = "Conclusion"
content3 = slide3.shapes.placeholders[0].text_frame
content3.text = "Summarize the key points from the presentation."

# Save the presentation to a file
prs.save("presentation_layout.pptx")
